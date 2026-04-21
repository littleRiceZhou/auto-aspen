import os
from pathlib import Path
import re
from typing import Optional

from pptx import Presentation
from pptx.util import Pt

# 导入现有的配置
try:
    from auto_aspen.docx_pdf import create_replacement_dict, sort_auto_aspen_keys_reverse, file_dir
except ImportError:
    # 备选方案，如果无法从 docx_pdf 导入
    file_dir = "static/re"
    def create_replacement_dict(parameter_values=None):
        return parameter_values or {}
    def sort_auto_aspen_keys_reverse(replacements):
        return sorted(replacements.keys(), key=len, reverse=True)


def _apply_replacements_to_text(full_text: str, sorted_keys: list, replacements: dict) -> str:
    """
    在整段字符串上应用替换，与 docx 中 auto_aspen 边界规则一致。
    长键优先（由 sort_auto_aspen_keys_reverse 保证）。
    """
    out = full_text
    for old_text in sorted_keys:
        if old_text not in out:
            continue
        new_text = str(replacements[old_text])
        if old_text.startswith("auto_aspen_"):
            escaped = re.escape(old_text)
            pattern = escaped + r"(?=\D|$)"
            out = re.sub(pattern, new_text, out)
        else:
            out = out.replace(old_text, new_text)
    return out


def _set_paragraph_font_size_pt(paragraph, size_pt: float) -> None:
    """将段落内所有 run 设为指定字号（磅）。"""
    if paragraph is None or size_pt is None or size_pt <= 0:
        return
    try:
        size = Pt(float(size_pt))
    except (TypeError, ValueError):
        return
    for run in paragraph.runs:
        run.font.size = size


def replace_text_in_shape(shape, replacements, font_size_pt: Optional[float] = 12.0):
    """
    在 PPT 形状中替换文本。

    PowerPoint 常把占位符拆成多个 run（例如 auto_aspen_2 | 0），按 run 替换会失败；
    必须在段落级别读取合并后的文本再写回（与 Word 跨 run 逻辑一致）。
    """
    replaced_count = 0

    text_frame = None
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        text_frame = shape.text_frame
    elif hasattr(shape, "text_frame"):
        text_frame = shape.text_frame

    if not text_frame:
        return 0

    sorted_keys = sort_auto_aspen_keys_reverse(replacements)

    for paragraph in text_frame.paragraphs:
        full = paragraph.text
        new_full = _apply_replacements_to_text(full, sorted_keys, replacements)
        if new_full != full:
            paragraph.text = new_full
            _set_paragraph_font_size_pt(paragraph, font_size_pt)
            replaced_count += 1

    return replaced_count

def replace_text_in_pptx(
    pptx_path,
    replacements,
    output_pptx_path=None,
    image_replacements=None,
    font_size_pt: Optional[float] = 12.0,
):
    """
    读取 pptx 文件，替换指定文本和图片并保存
    """
    print(f"正在读取 PPT 模板: {pptx_path}")
    
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"找不到模板文件: {pptx_path}")
        
    prs = Presentation(pptx_path)
    total_replaced = 0
    images_replaced = 0
    
    # 遍历所有幻灯片
    for slide in prs.slides:
        # 遍历幻灯片中的所有形状
        for shape in slide.shapes:
            # 处理图片替换 (通过占位符文本所在的形状位置插入图片)
            if image_replacements:
                for placeholder_text, img_info in image_replacements.items():
                    # 检查形状是否包含占位符文本
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if placeholder_text in shape.text_frame.text:
                            img_path = img_info.get("image_path")
                            if img_path and os.path.exists(img_path):
                                # 获取原形状的位置和大小
                                left = shape.left
                                top = shape.top
                                width = shape.width
                                height = shape.height
                                
                                # 在相同位置插入新图片
                                slide.shapes.add_picture(img_path, left, top, width, height)
                                
                                # 移除原来的占位符形状
                                sp = shape._element
                                sp.getparent().remove(sp)
                                
                                images_replaced += 1
                                print(f"成功在 PPT 中替换图片: {placeholder_text} -> {img_path}")
                                continue # 形状已被删除，跳过后续处理

            # 处理普通文本形状
            total_replaced += replace_text_in_shape(shape, replacements, font_size_pt)
            
            # 处理表格
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        total_replaced += replace_text_in_shape(cell, replacements, font_size_pt)
            
            # 处理组合形状
            if shape.shape_type == 6:  # Group shape
                for sub_shape in shape.shapes:
                    total_replaced += replace_text_in_shape(sub_shape, replacements, font_size_pt)

    # 确定输出路径
    if output_pptx_path is None:
        pptx_file = Path(pptx_path)
        output_pptx_path = os.path.join(file_dir, f"{pptx_file.stem}_generated.pptx")
    
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    prs.save(output_pptx_path)
    print(f"PPT 文档处理完成，共替换 {total_replaced} 处文本，{images_replaced} 处图片，已保存至: {output_pptx_path}")
    
    return output_pptx_path

def generate_pptx_document(
    parameters=None,
    text_to_images=None,
    output_name=None,
    template_path=None,
    font_size_pt: Optional[float] = 14.0,
):
    """
    封装的 PPT 生成函数，供外部调用。

    template_path: 可选，默认 models/production-template.pptx；EC 方案可用 models/ec-template.pptx。
    font_size_pt: 替换占位符后的正文字号（磅），默认 12，与 Word 技术方案一致；传 None 则不强制字号。
    """
    if template_path is None:
        template_path = "models/production-template.pptx"

    if not os.path.exists(template_path):
        return {
            "success": False,
            "error": f"PPT 模板文件不存在: {template_path}",
            "pptx_path": None
        }
    
    # 创建替换字典
    replacements = create_replacement_dict(parameters)
    
    # 确定输出文件名
    if output_name is None:
        output_name = "production_report_generated"
    
    output_pptx_path = f"{file_dir}/{output_name}.pptx"

    env_pt = os.getenv("AUTO_ASPEN_PPT_FONT_PT", "").strip()
    if env_pt:
        try:
            font_size_pt = float(env_pt)
        except ValueError:
            pass

    try:
        final_path = replace_text_in_pptx(
            template_path,
            replacements,
            output_pptx_path,
            image_replacements=text_to_images,
            font_size_pt=None,
        )
        return {
            "success": True,
            "pptx_path": final_path,
            "parameters_replaced": len(replacements),
            "images_replaced": 1 if text_to_images else 0,
        }
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "pptx_path": None
        }
